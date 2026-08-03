"""
Builds an ultra-professional, 6-slide presentation in PPTX format with:
- Dark Tech Navy Blue Background (RGB 13, 27, 42)
- High-contrast visual cards, badges, statistics counters, and timelines
- Embedded diagrams/images (from media__1784830302223.jpg and media__1784830308535.jpg)
- Comprehensive content tailored for presenting to professors:
  1. Title Slide: Project Overview & Academic Goals
  2. Problem Statement & Market Context (Citations: Grand View Research, Upstream Security, Fakhfakh 2022)
  3. Prior Trends & Literature Review (Timeline: 2018-2024, Research Gap identification)
  4. System Architecture & Dual-Core RTOS Flow (Embedded Architecture Diagram Image)
  5. Implementation Strategy, 12-Week Roadmap & Itemised Budget (Rs. 2,350 total)
  6. Expected Outcomes, Patentability & Academic Deliverables (System Overview Diagram Image)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─────────────────────────────────────────────────────────────
# CONSTANTS & PALETTE
# ─────────────────────────────────────────────────────────────
BG_COLOR     = RGBColor(13, 27, 42)     # Dark Tech Navy (#0D1B2A)
CARD_BG      = RGBColor(27, 38, 59)     # Deep Navy Card (#1B263B)
ACCENT_BLUE  = RGBColor(65, 90, 119)    # Accent Border (#415A77)
CYAN_GLOW    = RGBColor(0, 180, 216)    # High-tech Cyan (#00B4D8)
WHITE        = RGBColor(255, 255, 255)  # Text Pure White
MUTED_TEXT   = RGBColor(224, 225, 221)  # Muted White (#E0E1DD)
MUTED_GRAY   = RGBColor(141, 153, 174)  # Secondary Gray
GREEN_ACCENT = RGBColor(46, 196, 182)   # Success Green (#2EC4B6)
RED_ACCENT   = RGBColor(230, 57, 70)    # Alert Red (#E63946)

IMG_ARCH = r"C:\Users\mksin\.gemini\antigravity\brain\7a86b56c-8808-46db-bf91-4448eff62e7d\media__1784830302223.jpg"
IMG_FLOW = r"C:\Users\mksin\.gemini\antigravity\brain\7a86b56c-8808-46db-bf91-4448eff62e7d\media__1784830308535.jpg"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

def set_slide_background(slide):
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG_COLOR
    bg_shape.line.fill.background()
    return bg_shape

def add_header(slide, title_text, category_text):
    # Category / Badge
    tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_cat = tb_cat.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = CYAN_GLOW

    # Main Title
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=ACCENT_BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

# ─────────────────────────────────────────────────────────────
# SLIDE 1: TITLE SLIDE
# ─────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1)

# Subtle background shape accents
dec1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(-1.5), Inches(5.0), Inches(5.0))
dec1.fill.solid()
dec1.fill.fore_color.rgb = CARD_BG
dec1.line.fill.background()

# Title Badge
tb = s1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = "MINI PROJECT PROPOSAL | B.TECH EEE (2ND YEAR) | GCET GREATER NOIDA"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = CYAN_GLOW

# Main Title
tb = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.6))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "Cyber-Hardened Battery Management System\nfor Electric Vehicles"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# Subtitle
tb = s1.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.3), Inches(0.8))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "EKF-Based State Estimation with ML Intrusion Detection on CAN Bus"
p.font.size = Pt(18)
p.font.color.rgb = MUTED_TEXT

# Key Feature Pills
pills = [
    ("HARDWARE + SIMULATION", "Dual-Core ESP32 + BQ76920 AFE Prototype"),
    ("PATENT-READY CONCEPT", "Dynamic Noise Covariance R-Scaling"),
    ("IEEE PAPER TARGET", "IEEE APEC / ICIT Conference Submission")
]

for idx, (title, desc) in enumerate(pills):
    left = Inches(1.0 + idx * 3.8)
    card = add_card(s1, left, Inches(4.8), Inches(3.5), Inches(1.5))
    tb = s1.shapes.add_textbox(left + Inches(0.15), Inches(4.9), Inches(3.2), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_GLOW
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = MUTED_TEXT

# ─────────────────────────────────────────────────────────────
# SLIDE 2: PROBLEM STATEMENT & MARKET CONTEXT
# ─────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2)
add_header(s2, "Why This Project Matters: Problem Statement & Market Context", "Market Reality & Cyber Vulnerabilities")

# Stat Cards across the top
stats = [
    ("$54.4 Billion", "India EV Market (2025)", CYAN_GLOW),
    ("900+ Incidents", "Automotive Cyber Attacks (2021)", RED_ACCENT),
    ("19% CAGR", "India EV Growth Rate", GREEN_ACCENT),
    ("ZERO Security", "Built-in Protection on CAN Bus", RED_ACCENT)
]

for idx, (val, lbl, col) in enumerate(stats):
    left = Inches(0.8 + idx * 2.95)
    add_card(s2, left, Inches(1.6), Inches(2.75), Inches(1.3))
    tb = s2.shapes.add_textbox(left + Inches(0.1), Inches(1.7), Inches(2.55), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = val
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = col
    
    p2 = tf.add_paragraph()
    p2.text = lbl
    p2.font.size = Pt(10)
    p2.font.color.rgb = MUTED_TEXT

# Two Column Cards: The Problem vs Our Solution
# Problem Card (Left)
add_card(s2, Inches(0.8), Inches(3.1), Inches(5.7), Inches(3.6), border_color=RED_ACCENT)
tb = s2.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(5.3), Inches(3.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "THE PROBLEM: CAN BUS VULNERABILITY"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RED_ACCENT

prob_bullets = [
    "CAN Bus has NO authentication: Any compromised node or OBD-II dongle can inject malicious frames.",
    "BMS trusts incoming data blindly: Spoofed cell voltages or currents distort Kalman Filter state estimates.",
    "Severe SoC Estimation Errors: Corrupted data leads to >18% State-of-Charge drift during cyber-attacks.",
    "Safety Hazards: Risk of battery fire, sudden power loss, or premature cell degradation."
]
for b in prob_bullets:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED_TEXT
    p.space_before = Pt(4)

# Solution Card (Right)
add_card(s2, Inches(6.8), Inches(3.1), Inches(5.7), Inches(3.6), border_color=GREEN_ACCENT)
tb = s2.shapes.add_textbox(Inches(7.0), Inches(3.2), Inches(5.3), Inches(3.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "OUR SOLUTION: CYBER-HARDENED BMS"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = GREEN_ACCENT

sol_bullets = [
    "On-device ML Intrusion Detection: Random Forest model runs on ESP32 Core 0 to monitor bus traffic in real-time.",
    "Adaptive Covariance Scaling: Dynamically inflates EKF measurement noise R_eff = R_base × exp(10 × S_anomaly).",
    "Kalman Gain Suppression: Drives K → 0 under attack, forcing EKF to ignore false sensor data and rely on battery model.",
    "First Hardware Prototype: Combines ML-IDS + EKF feedback loop on low-cost dual-core MCU."
]
for b in sol_bullets:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED_TEXT
    p.space_before = Pt(4)

# Footer Citations
tb_cit = s2.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
tf_cit = tb_cit.text_frame
p_cit = tf_cit.paragraphs[0]
p_cit.text = "Citations: Grand View Research (India EV Market 2025); Upstream Security Automotive Cybersecurity Report (2022); Fakhfakh et al., Library Hi Tech (2022)."
p_cit.font.size = Pt(9)
p_cit.font.italic = True
p_cit.font.color.rgb = MUTED_GRAY

# ─────────────────────────────────────────────────────────────
# SLIDE 3: PRIOR TRENDS & LITERATURE REVIEW
# ─────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3)
add_header(s3, "Prior Trends & Literature Review: Identifying the Research Gap", "State of the Art vs. Our Innovation")

trends = [
    ("2018", "Seo et al. (IEEE PST)", "GAN-based IDS for CAN bus", "Purely software simulation; no micro-controller or BMS integration."),
    ("2020", "SAVIOR (USENIX Sec.)", "Physical invariants for AV security", "Focused on general Autonomous Vehicles, not specific to BMS cell dynamics."),
    ("2022", "Fakhfakh et al. (LHT)", "SLR: 30+ CAN attack vectors", "Comprehensive attack survey; no real-time embedded detection mechanism."),
    ("2023", "Perakovic et al. (MDPI)", "SVM/DT IDS on Kia Soul dataset", "High accuracy (99.9%), but computationally heavy and NOT deployed on MCU."),
    ("2024", "Kumar & Singh (PES)", "MCFO-DANN IDS for EV CAN", "Outperforms traditional ML but requires GPU infrastructure; no EKF coupling."),
    ("2024", "IEEE Xplore Trends", "EKF + Neural Network for SoC", "Focuses exclusively on estimation accuracy; zero cyber-attack awareness.")
]

# Grid 2x3 layout for timeline
for idx, (yr, auth, title, gap) in enumerate(trends):
    col = idx % 3
    row = idx // 3
    left = Inches(0.8 + col * 3.9)
    top = Inches(1.6 + row * 2.3)
    
    add_card(s3, left, top, Inches(3.7), Inches(2.1))
    tb = s3.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), Inches(3.5), Inches(1.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = f"{yr} | {auth}"
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_GLOW
    
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(2)
    
    p3 = tf.add_paragraph()
    p3.text = "Limitation: " + gap
    p3.font.size = Pt(10)
    p3.font.color.rgb = MUTED_TEXT
    p3.space_before = Pt(4)

# Research Gap Banner at bottom
add_card(s3, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), bg_color=CARD_BG, border_color=CYAN_GLOW)
tb_gap = s3.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.8))
tf_gap = tb_gap.text_frame
tf_gap.word_wrap = True
p_gap = tf_gap.paragraphs[0]
p_gap.text = "IDENTIFIED RESEARCH GAP:"
p_gap.font.size = Pt(11)
p_gap.font.bold = True
p_gap.font.color.rgb = CYAN_GLOW

p_gap2 = tf_gap.add_paragraph()
p_gap2.text = "No prior work couples an on-device CAN Intrusion Detection System directly into an EKF Measurement Noise Covariance Matrix (R) on a dual-core embedded MCU for BMS state protection."
p_gap2.font.size = Pt(11)
p_gap2.font.color.rgb = WHITE

# ─────────────────────────────────────────────────────────────
# SLIDE 4: SYSTEM ARCHITECTURE & DUAL-CORE RTOS FLOW
# ─────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
set_slide_background(s4)
add_header(s4, "System Architecture: Hardware Integration & Dual-Core RTOS Flow", "How Everything Works Together")

# Left Column: Image diagram container
add_card(s4, Inches(0.8), Inches(1.6), Inches(6.2), Inches(5.4))
if os.path.exists(IMG_ARCH):
    s4.shapes.add_picture(IMG_ARCH, Inches(0.9), Inches(1.7), width=Inches(6.0))

# Right Column: Bullet Points & Core Division
add_card(s4, Inches(7.2), Inches(1.6), Inches(5.3), Inches(5.4))
tb_arch = s4.shapes.add_textbox(Inches(7.4), Inches(1.7), Inches(4.9), Inches(5.2))
tf_arch = tb_arch.text_frame
tf_arch.word_wrap = True

p = tf_arch.paragraphs[0]
p.text = "DUAL-CORE FREERTOS ARCHITECTURE"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = CYAN_GLOW

arch_points = [
    ("Core 0 (Security Core):", "Runs high-priority CAN RX task. Extracts inter-arrival time (Δt) and frame frequency. Executes Random Forest C++ model (m2cgen) in <0.35 ms."),
    ("Core 1 (BMS Control Core):", "Reads cell voltages, current & temperature from TI BQ76920 AFE via I2C. Executes 1RC EKF state estimation and passive balancing."),
    ("Inter-Core Communication:", "Core 0 writes Anomaly Score S ∈ [0, 1] to non-blocking FreeRTOS queue (xQueueOverwrite). Core 1 reads S during EKF update cycle."),
    ("Dynamic R-Scaling Logic:", "R_eff = R_base × exp(10 × S). When S → 1.0 (attack), R_eff inflates 22,026×, suppressing Kalman gain K → 0."),
    ("Hardware Peripherals:", "Dual SN65HVD230 CAN transceivers (500 kbps), SSD1306 OLED display, MicroSD SPI logging, IRLML2502 MOSFET balancing.")
]

for title, body in arch_points:
    p1 = tf_arch.add_paragraph()
    p1.text = title
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_before = Pt(6)
    
    p2 = tf_arch.add_paragraph()
    p2.text = body
    p2.font.size = Pt(10)
    p2.font.color.rgb = MUTED_TEXT

# ─────────────────────────────────────────────────────────────
# SLIDE 5: IMPLEMENTATION STRATEGY, ROADMAP & BUDGET
# ─────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
set_slide_background(s5)
add_header(s5, "Implementation Roadmap & Budget: 12-Week Execution Plan", "Step-by-Step Methodology & Component Cost")

# Left Side: 12-Week Timeline
add_card(s5, Inches(0.8), Inches(1.6), Inches(6.8), Inches(5.4))
tb_t = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(6.4), Inches(5.1))
tf_t = tb_t.text_frame
tf_t.word_wrap = True

p = tf_t.paragraphs[0]
p.text = "12-WEEK STEP-BY-STEP ROADMAP"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = CYAN_GLOW

roadmap = [
    ("Weeks 1–3:", "Theoretical modeling, LTspice passive balancing circuit simulation & MATLAB/Simulink EKF algorithm setup."),
    ("Weeks 4–5:", "Hardware component procurement from Robu.in & ElectroPi.in. Safe bench wiring and power supply testing."),
    ("Weeks 6–7:", "ESP32 FreeRTOS firmware development: I2C driver for BQ76920, TWAI CAN driver, and basic EKF execution."),
    ("Weeks 8–9:", "Attacker ESP32 programming (DoS, Spoofing, Replay) & CAN dataset collection (20,000 frames) in CSV format."),
    ("Weeks 10–11:", "Random Forest IDS training in Python scikit-learn, export to C++ via m2cgen, and dual-core integration."),
    ("Week 12:", "System validation under active attack, IEEE conference paper writing & Indian Provisional Patent Form 2 filing.")
]

for w, d in roadmap:
    p1 = tf_t.add_paragraph()
    p1.text = w + " " + d
    p1.font.size = Pt(10.5)
    p1.font.color.rgb = MUTED_TEXT
    p1.space_before = Pt(4)

# Right Side: Itemised Budget Card
add_card(s5, Inches(7.8), Inches(1.6), Inches(4.7), Inches(5.4), border_color=GREEN_ACCENT)
tb_b = s5.shapes.add_textbox(Inches(8.0), Inches(1.7), Inches(4.3), Inches(5.1))
tf_b = tb_b.text_frame
tf_b.word_wrap = True

p = tf_b.paragraphs[0]
p.text = "ITEMISED PROJECT BUDGET (INR)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = GREEN_ACCENT

bom_items = [
    ("2x ESP32 Dev Boards", "Rs. 454"),
    ("TI BQ76920 AFE Breakout", "Rs. 1,000"),
    ("2x SN65HVD230 CAN Modules", "Rs. 160"),
    ("4x 18650 Li-ion Cells (1500mAh)", "Rs. 396"),
    ("4x IRLML2502 MOSFETs + 47Ω", "Rs. 120"),
    ("0.96\" SSD1306 OLED Display", "Rs. 145"),
    ("MicroSD Card Module + Misc PCB", "Rs. 76"),
    ("GRAND TOTAL", "Rs. 2,351")
]

for item, price in bom_items:
    p = tf_b.add_paragraph()
    p.text = f"{item}: {price}"
    p.font.size = Pt(11)
    if "GRAND TOTAL" in item:
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = GREEN_ACCENT
        p.space_before = Pt(10)
    else:
        p.font.color.rgb = MUTED_TEXT
        p.space_before = Pt(3)

p_note = tf_b.add_paragraph()
p_note.text = "\nPer Team Member (5 students): ~Rs. 470\nZero GPU Required | All Software Tools Free"
p_note.font.size = Pt(10)
p_note.font.italic = True
p_note.font.color.rgb = CYAN_GLOW

# ─────────────────────────────────────────────────────────────
# SLIDE 6: EXPECTED OUTCOMES, PATENTABILITY & CONCLUSION
# ─────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
set_slide_background(s6)
add_header(s6, "Expected Outcomes, Patentability & Academic Deliverables", "Why This Project Will Stand Out")

# Left Column: Image diagram container
add_card(s6, Inches(0.8), Inches(1.6), Inches(6.2), Inches(5.4))
if os.path.exists(IMG_FLOW):
    s6.shapes.add_picture(IMG_FLOW, Inches(0.9), Inches(1.7), width=Inches(6.0))

# Right Column: Patentability & Deliverables
add_card(s6, Inches(7.2), Inches(1.6), Inches(5.3), Inches(5.4), border_color=CYAN_GLOW)
tb_out = s6.shapes.add_textbox(Inches(7.4), Inches(1.7), Inches(4.9), Inches(5.2))
tf_out = tb_out.text_frame
tf_out.word_wrap = True

p = tf_out.paragraphs[0]
p.text = "TARGET DELIVERABLES"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = CYAN_GLOW

delivs = [
    ("1. Working Hardware Prototype:", "Demonstrating accurate SoC estimation under active 500 kbps CAN DoS and spoofing attacks (<1.4% SoC error)."),
    ("2. Indian Patent Application (Form 2):", "Covering system-level integration of CAN ML anomaly scoring with dynamic EKF measurement noise scaling."),
    ("3. IEEE Conference Paper:", "Targeting IEEE APEC / ICIT with experimental performance metrics, ROC-AUC curves, and latency benchmarks."),
    ("4. Why This Idea Stands Out:", "Exemplifies hands-on multi-disciplinary mastery (FreeRTOS, PCB design, ML deployment, embedded control) proving industry readiness.")
]

for title, desc in delivs:
    p1 = tf_out.add_paragraph()
    p1.text = title
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_before = Pt(6)
    
    p2 = tf_out.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = MUTED_TEXT

# Save PPTX
pptx_out = r"c:\Users\mksin\Desktop\AI hardened BMS\Cyber_Hardened_BMS_Proposal_Deck.pptx"
prs.save(pptx_out)
print(f"SUCCESSFULLY CREATED PPTX PRESENTATION AT: {pptx_out}")
