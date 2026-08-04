import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Image Asset Setup
ASSETS_DIR = r"c:\Users\mksin\Desktop\AI hardened BMS\assets"
os.makedirs(ASSETS_DIR, exist_ok=True)
SRC_IMAGE = r"C:\Users\mksin\.gemini\antigravity\brain\0c7977e6-c5f6-4522-8deb-7765a9c616f0\.user_uploaded\media_1785876165433.jpg"
BLUEPRINT_IMG_PATH = os.path.join(ASSETS_DIR, "blueprint_3d.png")
shutil.copy(SRC_IMAGE, BLUEPRINT_IMG_PATH)

# 2. Sleek High-Tech Dark Theme Palette
BG_COLOR      = RGBColor(10, 17, 30)      # Deep Dark Navy (#0A111E)
CARD_BG       = RGBColor(18, 28, 48)      # Card Background (#121C30)
CARD_BORDER   = RGBColor(30, 58, 95)      # Card Border (#1E3A5F)
CYAN_ACCENT   = RGBColor(0, 210, 255)     # Electric Cyan (#00D2FF)
GREEN_ACCENT  = RGBColor(16, 233, 122)    # Cyber Green (#10E97A)
WHITE         = RGBColor(255, 255, 255)   # Pure White
MUTED_TEXT    = RGBColor(180, 195, 215)   # Muted Blue-Gray (#B4C3D7)
RED_ALERT     = RGBColor(255, 75, 75)     # Alert Red (#FF4B4B)
YELLOW_WARN   = RGBColor(245, 158, 11)    # Amber Yellow (#F59E0B)

# 3. Helper Functions
def set_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text, category_text="CYBER-HARDENED BMS ARCHITECTURE"):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.13), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = CYAN_ACCENT
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.space_before = Pt(2)

def add_card(slide, left, top, width, height, border_color=CARD_BORDER, bg_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

# 4. Initialize 16:9 Widescreen Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: TITLE & PROBLEM STATEMENT VS SOLUTION
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1)

# Title Banner
add_card(s1, Inches(0.6), Inches(0.5), Inches(12.13), Inches(2.2), border_color=CYAN_ACCENT)
tb1 = s1.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.8))
tf1 = tb1.text_frame; tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "Self-Healing Cyber-Hardened Battery Management System (BMS)"
p1.font.size = Pt(25); p1.font.bold = True; p1.font.color.rgb = WHITE

p2 = tf1.add_paragraph()
p2.text = "Dual-Core Edge ML Intrusion Detection & Adaptive EKF Covariance Scaling for EV Safety"
p2.font.size = Pt(14); p2.font.color.rgb = CYAN_ACCENT; p2.space_before = Pt(4)

p3 = tf1.add_paragraph()
p3.text = "GCET EEE Department | Hardware Prototype Solution | Target Cost: ₹2,100 – ₹2,400 ($28 - $32 USD)"
p3.font.size = Pt(11); p3.font.color.rgb = MUTED_TEXT; p3.space_before = Pt(6)

# Section Header
tb_hdr1 = s1.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12.13), Inches(0.4))
tf_h1 = tb_hdr1.text_frame
ph1 = tf_h1.paragraphs[0]
ph1.text = "REAL-WORLD EV CYBERSECURITY CHALLENGES & PROPOSED HARDWARE SOLUTIONS"
ph1.font.size = Pt(12.5); ph1.font.bold = True; ph1.font.color.rgb = GREEN_ACCENT

# 4 Problem vs Solution Grid Cards
cards_data = [
    ("PROBLEM 1: Unauthenticated CAN/BLE Buses", "NO ENCRYPTION on EV 2W/3W buses allows attackers to flood DoS & spoof false telemetry frames, causing severe BMS state divergence.",
     "SOLUTION 1: Layer 1 & 2 Auth Engine", "ECDH SECP256R1 key exchange + AES-128-GCM BLE authentication & HMAC-SHA256 CAN payload signatures with 64-nonce replay cache.", RED_ALERT),
    
    ("PROBLEM 2: EKF Divergence Under Cyber Attack", "Conventional EKFs update state estimates on corrupted voltage/current inputs, driving State-of-Charge (SoC) errors >18.4%.",
     "SOLUTION 2: Layer 4 Adaptive EKF Scaling", "Dynamically scale covariance R_eff = R_base * exp(10*S). Under attack (S->1), K->0, suppressing corrupted sensor input while Coulomb counting continues.", CYAN_ACCENT),
    
    ("PROBLEM 3: Diagnostic Session Hijacking", "Malicious nodes exploit UDS (ISO 14229) OBD-II diagnostic requests (0x7E0) to release security seeds and force cell imbalances.",
     "SOLUTION 3: Layer 3 TinyML & UDS Inspector", "Core 0 inspector filters 0x27 SecurityAccess & 0x3E TesterPresent requests, overriding anomaly score S = 0.96 and blocking unauthorized writes.", GREEN_ACCENT),
    
    ("PROBLEM 4: High Latency & Hardware Cost", "Python ML frameworks require expensive GPUs and exceed real-time CAN frame budgets (>2.0 ms), freezing safety loops.",
     "SOLUTION 4: Zero-GPU C++ m2cgen Inference", "Random Forest decision trees compiled into native C++ IF/ELSE blocks running in <0.35 ms on ESP32 Core 0 (38.4 KB SRAM footprint).", YELLOW_WARN)
]

for idx, (p_title, p_desc, s_title, s_desc, col) in enumerate(cards_data):
    col_idx = idx % 2
    row_idx = idx // 2
    left = Inches(0.6 + col_idx * 6.16)
    top = Inches(3.4 + row_idx * 1.9)
    
    add_card(s1, left, top, Inches(5.97), Inches(1.75), border_color=col)
    tb_c = s1.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), Inches(5.61), Inches(1.5))
    tf_c = tb_c.text_frame; tf_c.word_wrap = True
    
    pt = tf_c.paragraphs[0]; pt.text = p_title; pt.font.size = Pt(11.5); pt.font.bold = True; pt.font.color.rgb = col
    pd = tf_c.add_paragraph(); pd.text = p_desc; pd.font.size = Pt(9.5); pd.font.color.rgb = MUTED_TEXT; pd.space_before = Pt(2)
    
    st = tf_c.add_paragraph(); st.text = s_title; st.font.size = Pt(11); st.font.bold = True; st.font.color.rgb = WHITE; st.space_before = Pt(4)
    sd = tf_c.add_paragraph(); sd.text = s_desc; sd.font.size = Pt(9.5); sd.font.color.rgb = MUTED_TEXT; sd.space_before = Pt(2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: SYSTEM BLUEPRINT ARCHITECTURE (IMAGE SLIDE)
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2)
add_header(s2, "3D Hardware Blueprint & Cyber-Hardened Architecture Diagram")

# Embed Blueprint Image (Left Container)
s2.shapes.add_picture(BLUEPRINT_IMG_PATH, Inches(0.6), Inches(1.35), width=Inches(8.5))

# Right Callouts Card
add_card(s2, Inches(9.3), Inches(1.35), Inches(3.43), Inches(5.65), border_color=CYAN_ACCENT)
tb_specs = s2.shapes.add_textbox(Inches(9.5), Inches(1.5), Inches(3.03), Inches(5.35))
tf_s = tb_specs.text_frame; tf_s.word_wrap = True

ps_hdr = tf_s.paragraphs[0]
ps_hdr.text = "HARDWARE CALLOUTS"; ps_hdr.font.size = Pt(13); ps_hdr.font.bold = True; ps_hdr.font.color.rgb = CYAN_ACCENT

callouts = [
    ("🔋 4S/16S Li-Ion Pack", "14.8V nominal 4S modular sub-unit (daisy-chainable up to 16S 57.6V)"),
    ("🛡️ TI BQ76920 AFE IC", "Analog Front-End monitoring 4S cell voltages & shunt current via I2C (0x18)"),
    ("🧠 ESP32 Master MCU", "Dual-Core 240MHz (Core 0 Security ML / Core 1 Deterministic EKF)"),
    ("⚡ High-Side SSR Cutoff", "GPIO 17 Optocoupler P-FET gate driver for pack isolation (<1.2 ms trip)"),
    ("⚡ Dual CAN Transceivers", "SN65HVD230 high-speed CAN 2.0B transceivers (500 kbps, 120Ω split termination)"),
    ("🔴 ESP32 Attacker Node", "Injection node triggering DoS, Spoofing, Replay, Fuzzing & UDS attacks"),
    ("☁️ TCU Telematics Gateway", "Streams real-time CAN telemetry (0x180) to Node-RED / Grafana dashboards")
]

for title, desc in callouts:
    p_t = tf_s.add_paragraph()
    p_t.text = title; p_t.font.size = Pt(10.5); p_t.font.bold = True; p_t.font.color.rgb = GREEN_ACCENT; p_t.space_before = Pt(5)
    p_d = tf_s.add_paragraph()
    p_d.text = desc; p_d.font.size = Pt(9); p_d.font.color.rgb = MUTED_TEXT

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: IMPLEMENTATION ARCHITECTURE & FREERTOS DUAL-CORE WORKING
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3)
add_header(s3, "Hardware Implementation & FreeRTOS Dual-Core Execution Model")

# Left Card: Dual-Core FreeRTOS Split
add_card(s3, Inches(0.6), Inches(1.35), Inches(5.9), Inches(5.65), border_color=CYAN_ACCENT)
tb_core = s3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.35))
tf_core = tb_core.text_frame; tf_core.word_wrap = True

pc_hdr = tf_core.paragraphs[0]
pc_hdr.text = "FREERTOS DUAL-CORE TASK SEPARATION"; pc_hdr.font.size = Pt(13); pc_hdr.font.bold = True; pc_hdr.font.color.rgb = CYAN_ACCENT

cores_detail = [
    ("Core 0: Security & AI Engine (Priority 5, 10ms)",
     "• CAN Ring Buffer Feature Extractor: Calculates inter-arrival dt, frequency f, DLC & payload entropy.\n"
     "• Zero-GPU TinyML Random Forest: m2cgen compiled C++ decision tree (ids_model.h) computes Anomaly Score S in <0.35 ms.\n"
     "• UDS (ISO 14229) Inspector: Intercepts CAN ID 0x7E0 SecurityAccess (0x27) requests, overriding S = 0.96.\n"
     "• Wireless BLE Stack: ECDH SECP256R1 key exchange & AES-128 telemetry encryption.", GREEN_ACCENT),
    
    ("Core 1: Battery Control & Adaptive EKF (Priority 4, 100ms)",
     "• AFE Acquisition: Polls TI BQ76920 cell voltages (VC1-VC4) & pack shunt current every 100 ms.\n"
     "• 2RC Thevenin ECM EKF: Propagates state x_k = [SoC, V1, V2]^T using Coulomb counting & RC dynamics.\n"
     "• Adaptive Covariance Scaling: R_eff = R_base * exp(10*S). Under attack (S->1), R_eff = 220.26 => Kalman Gain K->0.\n"
     "• Actuation & Self-Healing: Controls IRLML2502 passive balancing MOSFETs & GPIO 17 SSR cutoff.", CYAN_ACCENT)
]

for title, desc, col in cores_detail:
    pt = tf_core.add_paragraph()
    pt.text = title; pt.font.size = Pt(11); pt.font.bold = True; pt.font.color.rgb = col; pt.space_before = Pt(8)
    pd = tf_core.add_paragraph()
    pd.text = desc; pd.font.size = Pt(9.5); pd.font.color.rgb = MUTED_TEXT; pd.space_before = Pt(3)

# Right Card: Low-Cost BOM Table
add_card(s3, Inches(6.8), Inches(1.35), Inches(5.93), Inches(5.65), border_color=GREEN_ACCENT)
tb_bom_hdr = s3.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4))
tf_bh = tb_bom_hdr.text_frame
pbh = tf_bh.paragraphs[0]
pbh.text = "LOW-COST BOM MATRIX (TARGET ₹2,100 – ₹2,400)"; pbh.font.size = Pt(12.5); pbh.font.bold = True; pbh.font.color.rgb = GREEN_ACCENT

table_shape = s3.shapes.add_table(7, 3, Inches(7.0), Inches(2.05), Inches(5.53), Inches(4.7))
table = table_shape.table
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(2.13)
table.columns[2].width = Inches(1.2)

bom_data = [
    ("Component Category", "Optimized Choice", "Price (INR)"),
    ("Microcontrollers", "1x ESP32-S3 / 2x Bare ICs", "₹690"),
    ("CAN Transceivers", "2x VP230 / TJA1051 ICs", "₹240"),
    ("Battery Pack", "4S1P Reclaimed 18650 Pack", "₹450"),
    ("SSR Cutoff Driver", "Optocoupler + P-FET Switch", "₹90"),
    ("AFE & Sensors", "TI BQ76920 IC + Shunt", "₹800"),
    ("TOTAL BENCHMARK", "Complete Hardware Setup", "₹2,100 – ₹2,400")
]

for row_idx, row in enumerate(bom_data):
    for col_idx, text in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if col_idx < 2 else PP_ALIGN.RIGHT
        p.font.size = Pt(9.5)
        if row_idx == 0:
            p.font.bold = True; p.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(30, 58, 95)
        elif row_idx == 6:
            p.font.bold = True; p.font.color.rgb = GREEN_ACCENT
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(20, 40, 65)
        else:
            p.font.color.rgb = MUTED_TEXT
            cell.fill.solid(); cell.fill.fore_color.rgb = CARD_BG

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: DETAILED WORKING OF ALL 7 SECURITY & SELF-HEALING LAYERS
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4)
add_header(s4, "Detailed Working of All 7 Security & Self-Healing Defense Layers")

layers_working = [
    ("LAYER 1: BLE Peripheral Authentication",
     "Uses ECDH SECP256R1 key exchange to derive session keys, AES-128-GCM payload encryption, and a 64-nonce sliding cache to block replay attacks over wireless interfaces.", CYAN_ACCENT),
    
    ("LAYER 2: Bus Command HMAC Verification",
     "Computes HMAC-SHA256 frame signatures on incoming CAN write commands. Invalid payload signatures or stale nonces are dropped before reaching actuation loops.", GREEN_ACCENT),
    
    ("LAYER 3: TinyML IDS & UDS (ISO 14229) Inspector",
     "Evaluates 10-tree Random Forest exported via m2cgen into C++ IF/ELSE blocks (<0.35 ms) on Core 0. Intercepts CAN ID 0x7E0 SecurityAccess (0x27) requests, yielding Anomaly Score S in [0, 1].", CYAN_ACCENT),
    
    ("LAYER 4: Adaptive EKF Trust Covariance Scaling",
     "Dynamically scales measurement noise covariance R_eff = R_base * exp(10*S). Under cyber-attack (S->1), R_eff = 220.26 => Kalman Gain K->0. Suppresses corrupted sensor input y_k so x_k stays on model.", GREEN_ACCENT),
    
    ("LAYER 5: Self-Healing Decision Engine",
     "5-Stage Action Threshold State Machine:\n"
     "• S < 0.30: Action 0 (Normal)  |  • 0.30 <= S < 0.50: Action 1 (Reduce Charge Current via PWM)\n"
     "• 0.50 <= S < 0.70: Action 2 (Disable Cell Balancing)  |  • 0.70 <= S < 0.90: Action 3 (Limp Home 50% Throttle)\n"
     "• S >= 0.90: Action 4 (High-Side SSR Pack Isolation Cutoff on GPIO 17)", CYAN_ACCENT),
    
    ("LAYER 6: High-Side SSR Physical Pack Isolation",
     "Drives GPIO 17 HIGH to activate optocoupler PC817 + P-FET gate driver in <1.2 ms during severe attack (S >= 0.90), physically disconnecting the Li-ion battery pack from load bus.", RED_ALERT),
    
    ("LAYER 7: Secure Cloud TCU Telematics Gateway",
     "Streams encrypted telemetry over CAN ID 0x180 (SoC, Voltage, Current, Anomaly Score S, Action ID) to TCU node for secure MQTT cloud publishing to Grafana / Node-RED dashboards.", GREEN_ACCENT)
]

for idx, (title, desc, col) in enumerate(layers_working):
    top = Inches(1.35 + idx * 0.83)
    add_card(s4, Inches(0.6), top, Inches(12.13), Inches(0.76), border_color=col)
    tb_l = s4.shapes.add_textbox(Inches(0.8), top + Inches(0.06), Inches(11.7), Inches(0.64))
    tf_l = tb_l.text_frame; tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    r1 = p.add_run(); r1.text = title + "\n"; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = col
    r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(9.5); r2.font.color.rgb = MUTED_TEXT

# Save Presentation
PPTX_PATH = r"c:\Users\mksin\Desktop\AI hardened BMS\Cyber_Hardened_BMS_4Slide_Technical_Presentation.pptx"
prs.save(PPTX_PATH)
print(f"SUCCESSFULLY GENERATED TECHNICAL 4-SLIDE PPTX AT: {PPTX_PATH}")
